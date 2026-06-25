#!/usr/bin/env python3
"""
재사용 템플릿 — python-pptx로 편집가능 아키텍처 다이어그램(단일 슬라이드, 네이티브 도형).
사용법: NODES / EDGES / OUT / TITLE 만 수정 후 `python3 build_pptx_template.py`.
  - shape: 'rrect'(서비스) | 'db'(실린더/DB) | 'cube'(HW/GPU) | 'cloud'(외부 플랫폼)
  - 화살표는 parent.right-center -> child.left-center 자동 라우팅(좌->우 트리 기준)
채워진 예시는 NODES/EDGES/OUT/TITLE 를 대상 시스템 값으로 교체해 사용.
검증: 빌드 후 재오픈해 슬라이드/커넥터(p:cxnSp)/화살촉(a:tailEnd)/경계 확인. soffice 無 → 픽셀렌더 Unknown.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

FONT = '맑은 고딕'
OUT = 'output/architecture/diagram.pptx'
TITLE = '다이어그램 제목'

SHAPES = {'rrect': MSO_SHAPE.ROUNDED_RECTANGLE, 'db': MSO_SHAPE.CAN,
          'cube': MSO_SHAPE.CUBE, 'cloud': MSO_SHAPE.CLOUD}

# 색상 팔레트 예시: 서비스 파랑/보라/주황, DB 초록, 보안 분홍, HW 회색, 플랫폼 주황
# key: (label, sub, x, y, w, h, fill_hex, shape)   — 단위 inch, 슬라이드 16:9 = 13.333 x 7.5
NODES = {
    'FE': ('Frontend', ':8080', 0.4, 2.9, 2.1, 0.85, 'E1F5FF', 'rrect'),
    'BE': ('Backend', ':8000', 3.1, 2.9, 2.2, 0.85, 'F3E5F5', 'rrect'),
    'DB': ('Database', ':5432', 6.0, 2.9, 2.1, 0.85, 'E8F5E9', 'db'),
}
EDGES = [('FE', 'BE'), ('BE', 'DB')]  # (parent, child)


def set_run(r, text, size, bold=False, color='222222'):
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color)
    r.font.name = FONT
    rPr = r._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):  # ea 필수 — 누락 시 한글 tofu 위험
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', FONT)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.6))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), TITLE, 24, True, '1A1A1A')

rects = {}
for key, (label, sub, x, y, w, h, fill, shp) in NODES.items():
    sp = slide.shapes.add_shape(SHAPES[shp], Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    sp.line.color.rgb = RGBColor.from_string('9E9E9E'); sp.line.width = Pt(1)
    try:
        sp.shadow.inherit = False
    except Exception:
        pass
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    set_run(pp.add_run(), label, 12.5, True, '1A1A1A')
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        set_run(p2.add_run(), sub, 9, False, '555555')
    rects[key] = (x, y, w, h)


def arrow(x1, y1, x2, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = RGBColor.from_string('5A5A5A'); c.line.width = Pt(1.5)
    ln = c._element.spPr.get_or_add_ln()  # 화살촉: ln 얻어 a:tailEnd append
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))


for a, b in EDGES:  # parent.right-center -> child.left-center (좌->우)
    ax, ay, aw, ah = rects[a]; bx, by, bw, bh = rects[b]
    arrow(ax + aw, ay + ah / 2, bx, by + bh / 2)

prs.save(OUT)
print('saved', OUT)
