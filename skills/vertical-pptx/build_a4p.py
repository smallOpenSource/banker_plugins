#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
vertical-pptx — A4 세로(210 x 297mm) PPTX python 빌더.

build_a4p.js 와 의미가 같은 산출물을 낸다. 같은 deck JSON 을 먹이면 sldSz·슬라이드 수·
도형 x/y/cx/cy·텍스트·3슬롯 폰트·a:bodyPr·warnings 가 전부 일치해야 한다(교차 동등성 테스트가 본다).

★ 입력은 deck JSON 뿐이다. 마크다운 플래그는 여기 없다.
  마크다운 → deck JSON 변환은 node 빌더가 단독으로 맡는다. 파서가 두 벌이 되는 순간
  같은 원고가 두 개의 다른 덱이 되고, 그것이 이 설계가 없애려던 바로 그 드리프트다.

★ INV-G (기하 불변식) — node 빌더와 공유하는 절대 규칙.
  모든 도형의 x·y·cx·cy 는 다음 값들만의 순수 함수다:
    meta.marginMm · 페이지 인덱스 · 블록 순서 · 블록 type · bullets.items 배열 길이 ·
    명시된 rows · 컬럼 스팬.
  텍스트의 내용·길이·폰트·pt 는 어떤 좌표에도 들어가지 않는다.

★ 측정은 여기서 하지 않는다.
  폰트 바이너리 파싱·폭 합산·오버플로 판정은 verify_a4p.mjs 단독 소유다.
  이 파일은 폰트 "이름" 을 고르고 그 파일이 존재하는지만 확인해서 경로를 보고한다.

python-pptx 특유의 함정 셋. 전부 이 저장소에서 실측한 것이고, 안 고치면 조용히 틀린 파일이 나온다.
  F3  slide_width/height 만 바꾸면 <p:sldSz ... type="screen4x3"/> 가 남는다.
      치수는 A4 인데 선언은 4:3 인 파일이고, 그 선언을 믿는 뷰어·변환기가 있다 → type 속성을 지운다.
  F4  run.font.name 은 a:latin 하나만 낸다. a:ea / a:cs 가 없으면 한글이 대체 폰트로 떨어진다
      → 세 슬롯에 같은 typeface 를 직접 주입한다(패턴 출처: skills/arch-diagram/build_pptx_template.py).
  F13 add_textbox 의 기본 a:bodyPr 는 wrap="none" + <a:spAutoFit/> 다.
      텍스트가 안 감기고, 뷰어가 열 때 도형 크기를 바꾼다 — 파일 안 좌표는 맞는데 화면에서 INV-G 가 깨진다
      → wrap="square" · 인셋 4개 0 · rtlCol="0" · anchor="t" · autofit 요소 제거로 고정한다.

기하 (단위 EMU, 1mm = 36000, m = meta.marginMm):
  슬라이드 7560000 x 10692000 (고정)
  여백 36000m · 본문 (7560000 - 72000m) x (10692000 - 72000m) · 거터 144000 (4mm 고정)
  8열: 컬럼 폭 (182 - 2m) * 4500 · 20행: 행 높이 (221 - 2m) * 1800
  컬럼 c 에서 k 열: x = 36000m + (컬럼폭 + 144000) * c · cx = 컬럼폭 * k + 144000 * (k - 1)
  행 r 에서 n 행:   y = 36000m + (행높이 + 144000) * r · cy = 행높이 * n + 144000 * (n - 1)
  m=10 검산: 컬럼 폭 729000 · 행 높이 361800 · 피치 873000 / 505800

행 수 함수 (측정 없음):
  rows(block) = min(18, block.rows ?? default(type))
  default: head=1 · body=2 · bullets=max(1, items 개수) · 페이지 제목 = 2행 고정
  ★ 클램프 발동 조건은 block.rows 가 아니라 block.rows ?? default(type) 다.
    항목 19개짜리 불릿에 rows 를 안 적어도 클램프가 걸리고, 그때 경고가 없으면
    항목 17·18 이 조용히 겹친 채 출하된다.

페이지 채우기(세는 규칙만): avail = 20 - 2 = 18. used + rows(block) > 18 이면 새 슬라이드.
  블록은 쪼개지 않는다(쪼개려면 측정이 필요하다).

요구 사항: python 3.9 이상 + python-pptx.
  이 파일은 셰방에 의존하지 않는다 — 호출자가 인터프리터를 고른다.
  `python3` 가 3.9 미만인 머신이 있으므로(이 개발 머신의 python3 는 3.6.8 이다)
  실제로 쓸 인터프리터를 명시해 부르는 것이 안전하다: `/path/to/python build_a4p.py ...`.
  3.9 미만이면 exit 2 로 그 사실만 알린다. 의존성을 자동 설치하지 않는다.

종료 코드: 0 통과 · 1 검증 실패 · 2 환경 미비 · 3 사용법 오류.
"""

import json
import math
import os
import subprocess
import sys

# ── 상수 (build_a4p.js 와 값이 같아야 한다) ─────────────────────────────────
MIN_PYTHON = (3, 9)
EMU_PER_MM = 36000
SLIDE_W = 210 * EMU_PER_MM  # 7560000
SLIDE_H = 297 * EMU_PER_MM  # 10692000
GUTTER = 4 * EMU_PER_MM     #  144000
COLS = 8
ROWS = 20
TITLE_ROWS = 2
MAX_BLOCK_ROWS = ROWS - TITLE_ROWS  # 18
MARGIN_MM_MIN = 10  # 동결 스펙의 "여백 10mm 이상"
MARGIN_MM_MAX = 40
DEFAULT_MARGIN_MM = 10
DEFAULT_BODY_PT = 10.5  # 동결 범위 10~11pt
LINE_SPACING = 1.2
BULLET_PREFIX = '\u2022 '
COLOR_HEAD = '1A1A1A'
COLOR_BODY = '222222'
BLOCK_TYPES = ('head', 'body', 'bullets')
ORIENTATIONS = ('portrait', 'landscape')

EXIT_OK = 0
EXIT_VERIFY = 1
EXIT_ENV = 2
EXIT_USAGE = 3

# 폰트 체인. --font → BANKER_PPTX_FONT → meta.font → 아래 목록 → 시스템 한국어 sans.
FONT_CHAIN = ('Noto Sans CJK KR', 'Malgun Gothic', 'Apple SD Gothic Neo', 'NanumGothic')
# fc-match 가 없는 환경에서 파일명으로 찾기 위한 힌트(정규화된 소문자 조각).
FONT_FILE_HINTS = {
    'noto sans cjk kr': ('notosanscjkkr', 'notosanscjk', 'notosanskr'),
    'malgun gothic': ('malgun',),
    'apple sd gothic neo': ('applesdgothicneo', 'applegothic'),
    'nanumgothic': ('nanumgothic',),
}
FONT_EXT = frozenset(('.ttf', '.ttc', '.otf', '.otc'))


# ── 오류 타입 ───────────────────────────────────────────────────────────────
class UsageError(Exception):
    exit_code = EXIT_USAGE


class EnvError(Exception):
    exit_code = EXIT_ENV


class VerifyError(Exception):
    exit_code = EXIT_VERIFY


# ── 그리드 ──────────────────────────────────────────────────────────────────
def grid_for(margin_mm):
    m = margin_mm
    col_w = (182 - 2 * m) * 4500
    row_h = (221 - 2 * m) * 1800
    return {
        'marginMm': m,
        'margin': EMU_PER_MM * m,
        'slideW': SLIDE_W,
        'slideH': SLIDE_H,
        'contentW': SLIDE_W - 2 * EMU_PER_MM * m,
        'contentH': SLIDE_H - 2 * EMU_PER_MM * m,
        'gutter': GUTTER,
        'cols': COLS,
        'rows': ROWS,
        'colW': col_w,
        'rowH': row_h,
        'colPitch': col_w + GUTTER,
        'rowPitch': row_h + GUTTER,
    }


def col_x(g, c):
    return g['margin'] + g['colPitch'] * c


def col_cx(g, k):
    return g['colW'] * k + GUTTER * (k - 1)


def row_y(g, r):
    return g['margin'] + g['rowPitch'] * r


def row_cy(g, n):
    return g['rowH'] * n + GUTTER * (n - 1)


def default_rows(block):
    if block.get('type') == 'head':
        return 1
    if block.get('type') == 'body':
        return 2
    return max(1, len(block.get('items') or []))  # bullets


def requested_rows(block):
    """클램프 이전의 요청 행 수. 클램프 발동 판정은 반드시 이 값으로 한다."""
    rows = block.get('rows')
    return default_rows(block) if rows is None else rows


def rows_of(block):
    return min(MAX_BLOCK_ROWS, requested_rows(block))


def bullet_spans(n, count):
    """
    불릿 항목별 행 배분.
    항목 i 는 min(i, n-1) 행에서 시작하고, 남는 행은 전부 마지막 항목이 가져간다.
    n < 항목 수면 n-1 행부터 겹쳐 쌓인다 — 그 겹침을 보고하는 것은 검증기 몫이다.
    """
    spans = []
    for i in range(count):
        start = min(i, n - 1)
        end = n - 1 if i == count - 1 else min(i + 1, n - 1) - 1
        spans.append({'start': start, 'rows': max(1, end - start + 1)})
    return spans


# ── deck 검증 (전부 exit 3) ─────────────────────────────────────────────────
def is_plain_object(v):
    return isinstance(v, dict)


def as_integer(value):
    """
    JS 의 Number.isInteger 와 같은 판정.
    JSON 의 10.0 은 JS 에서 정수 10 이지만 python 에서는 float 이다 —
    그대로 isinstance(int) 로 보면 node 빌더가 받는 덱을 python 빌더만 거부한다.
    """
    if isinstance(value, bool):
        return None
    if isinstance(value, int):
        return value
    if isinstance(value, float) and math.isfinite(value) and value.is_integer():
        return int(value)
    return None


def as_number(value):
    if isinstance(value, bool):
        return None
    if isinstance(value, (int, float)) and math.isfinite(value):
        return float(value)
    return None


def show(value):
    """오류 메시지에 값을 그대로 되비추기 위한 JSON 표현(JS 의 JSON.stringify 와 같은 자리)."""
    try:
        return json.dumps(value, ensure_ascii=False)
    except (TypeError, ValueError):
        return repr(value)


def orientation_of(deck, page):
    meta = deck.get('meta') or {}
    return page.get('orientation') or meta.get('orientation') or 'portrait'


def validate_deck(deck):
    if not is_plain_object(deck):
        raise UsageError('deck JSON 최상위는 객체여야 한다.')
    meta = {} if deck.get('meta') is None else deck.get('meta')
    if not is_plain_object(meta):
        raise UsageError('meta 는 객체여야 한다.')

    if meta.get('marginMm') is not None:
        m = as_integer(meta.get('marginMm'))
        if m is None or m < MARGIN_MM_MIN or m > MARGIN_MM_MAX:
            raise UsageError(
                'meta.marginMm 은 {0} 이상 {1} 이하의 정수여야 한다 (받은 값: {2}).'.format(
                    MARGIN_MM_MIN, MARGIN_MM_MAX, show(meta.get('marginMm'))))
    if meta.get('bodyPt') is not None:
        pt = as_number(meta.get('bodyPt'))
        if pt is None or pt <= 0 or pt > 400:
            raise UsageError('meta.bodyPt 는 0 초과 400 이하의 수여야 한다 (받은 값: {0}).'.format(
                show(meta.get('bodyPt'))))
    if meta.get('title') is not None and not isinstance(meta.get('title'), str):
        raise UsageError('meta.title 은 문자열이어야 한다.')
    if meta.get('font') is not None and not isinstance(meta.get('font'), str):
        raise UsageError('meta.font 는 문자열이어야 한다.')
    if meta.get('orientation') is not None and meta.get('orientation') not in ORIENTATIONS:
        raise UsageError('meta.orientation 은 {0} 여야 한다.'.format(' 또는 '.join(ORIENTATIONS)))

    pages = deck.get('pages')
    if not isinstance(pages, list) or not pages:
        raise UsageError('pages 는 비어 있지 않은 배열이어야 한다.')

    for pi, page in enumerate(pages):
        at = 'pages[{0}]'.format(pi)
        if not is_plain_object(page):
            raise UsageError('{0} 는 객체여야 한다.'.format(at))
        if page.get('title') is not None and not isinstance(page.get('title'), str):
            raise UsageError('{0}.title 은 문자열이어야 한다.'.format(at))
        if page.get('orientation') is not None and page.get('orientation') not in ORIENTATIONS:
            raise UsageError('{0}.orientation 은 {1} 여야 한다.'.format(at, ' 또는 '.join(ORIENTATIONS)))
        blocks = [] if page.get('blocks') is None else page.get('blocks')
        if not isinstance(blocks, list):
            raise UsageError('{0}.blocks 는 배열이어야 한다.'.format(at))
        for bi, block in enumerate(blocks):
            bat = '{0}.blocks[{1}]'.format(at, bi)
            if not is_plain_object(block):
                raise UsageError('{0} 는 객체여야 한다.'.format(bat))
            if block.get('type') not in BLOCK_TYPES:
                raise UsageError('{0}.type 은 {1} 중 하나여야 한다 (받은 값: {2}).'.format(
                    bat, ' / '.join(BLOCK_TYPES), show(block.get('type'))))
            if block.get('type') == 'bullets':
                items = block.get('items')
                if not isinstance(items, list) or not items:
                    raise UsageError('{0}.items 는 비어 있지 않은 배열이어야 한다.'.format(bat))
                for ii, item in enumerate(items):
                    if not isinstance(item, str):
                        raise UsageError('{0}.items[{1}] 는 문자열이어야 한다.'.format(bat, ii))
            elif not isinstance(block.get('text'), str):
                raise UsageError('{0}.text 는 문자열이어야 한다.'.format(bat))
            if block.get('rows') is not None:
                rows = as_integer(block.get('rows'))
                if rows is None or rows < 1:
                    raise UsageError('{0}.rows 는 1 이상의 정수여야 한다 (받은 값: {1}).'.format(
                        bat, show(block.get('rows'))))

    # 방향 혼합은 파일 형식이 허용하지 않는다 — 하나의 pptx 는 하나의 sldSz 만 갖는다.
    seen = set(orientation_of(deck, p) for p in pages)
    if len(seen) > 1:
        raise UsageError(
            '한 pptx 는 세로와 가로를 함께 담을 수 없다(슬라이드 크기가 파일당 하나다). '
            '방향별로 파일을 나누어 각각 빌드해라 — split the deck into separate files, '
            'one per orientation.')
    if 'landscape' in seen:
        raise UsageError('이 스킬은 A4 세로(210x297mm) 전용이다. 가로 덱은 대상이 아니다.')
    return True


# ── layout: deck JSON → 좌표 (INV-G 의 구현체) ──────────────────────────────
def layout(deck):
    """
    반환: {'grid', 'bodyPt', 'slides', 'warnings'}
      slides[i]['shapes'][j] 는 role·pageIndex·blockIndex·itemIndex·rows·x·y·cx·cy·text·pt·bold·color.
      pageIndex·blockIndex 가 있어야 검증기가 (슬라이드, 도형) 을 (페이지, 블록) 으로 되돌릴 수 있다.
    """
    validate_deck(deck)
    meta = deck.get('meta') or {}
    margin_mm = DEFAULT_MARGIN_MM if meta.get('marginMm') is None else as_integer(meta.get('marginMm'))
    grid = grid_for(margin_mm)
    body_pt = DEFAULT_BODY_PT if meta.get('bodyPt') is None else as_number(meta.get('bodyPt'))
    pt = {'title': body_pt + 6, 'head': body_pt + 2, 'body': body_pt, 'bullets': body_pt}
    full_cx = col_cx(grid, COLS)
    slides = []
    warnings = []

    for page_index, page in enumerate(deck['pages']):
        title = '' if page.get('title') is None else page.get('title')
        state = {'slide': None, 'used': 0, 'part': 0}

        def open_slide():
            slide = {'index': len(slides), 'pageIndex': page_index, 'part': state['part'],
                     'title': title, 'shapes': []}
            state['part'] += 1
            state['used'] = 0
            slide['shapes'].append({
                'role': 'title',
                'pageIndex': page_index,
                'blockIndex': None,
                'itemIndex': None,
                'rows': TITLE_ROWS,
                'x': col_x(grid, 0),
                'y': row_y(grid, 0),
                'cx': full_cx,
                'cy': row_cy(grid, TITLE_ROWS),
                'text': title,
                'pt': pt['title'],
                'bold': True,
                'color': COLOR_HEAD,
            })
            slides.append(slide)
            state['slide'] = slide

        open_slide()

        blocks = [] if page.get('blocks') is None else page.get('blocks')
        for block_index, block in enumerate(blocks):
            want = requested_rows(block)
            n = rows_of(block)
            if want > MAX_BLOCK_ROWS:
                warnings.append({
                    'code': 'BLOCK_TOO_TALL',
                    'page': page_index,
                    'block': block_index,
                    'detail': 'rows {0} -> {1}'.format(want, n),
                })
            # 세는 규칙만: 넘치면 새 슬라이드로 넘긴다. 블록은 절대 쪼개지 않는다.
            if state['used'] > 0 and state['used'] + n > MAX_BLOCK_ROWS:
                open_slide()
            base = TITLE_ROWS + state['used']
            slide = state['slide']

            if block['type'] == 'bullets':
                items = block['items']
                for item_index, span in enumerate(bullet_spans(n, len(items))):
                    slide['shapes'].append({
                        'role': 'bullet',
                        'pageIndex': page_index,
                        'blockIndex': block_index,
                        'itemIndex': item_index,
                        'rows': span['rows'],
                        'x': col_x(grid, 0),
                        'y': row_y(grid, base + span['start']),
                        'cx': full_cx,
                        'cy': row_cy(grid, span['rows']),
                        'text': BULLET_PREFIX + items[item_index],
                        'pt': pt['bullets'],
                        'bold': False,
                        'color': COLOR_BODY,
                    })
            else:
                is_head = block['type'] == 'head'
                slide['shapes'].append({
                    'role': block['type'],
                    'pageIndex': page_index,
                    'blockIndex': block_index,
                    'itemIndex': None,
                    'rows': n,
                    'x': col_x(grid, 0),
                    'y': row_y(grid, base),
                    'cx': full_cx,
                    'cy': row_cy(grid, n),
                    'text': block['text'],
                    'pt': pt['head'] if is_head else pt['body'],
                    'bold': is_head,
                    'color': COLOR_HEAD if is_head else COLOR_BODY,
                })
            state['used'] += n

    return {'grid': grid, 'bodyPt': body_pt, 'slides': slides, 'warnings': warnings}


def assert_inside_slide(laid):
    """기하가 실제로 여백 안에 있는지 자기 점검. 클램프가 살아 있으면 절대 발동하지 않는다."""
    g = laid['grid']
    bad = []
    for s in laid['slides']:
        for j, sh in enumerate(s['shapes']):
            if (sh['x'] < g['margin'] or sh['y'] < g['margin']
                    or sh['x'] + sh['cx'] > SLIDE_W - g['margin']
                    or sh['y'] + sh['cy'] > SLIDE_H - g['margin']):
                bad.append('slide {0} shape {1} ({2}): x={3} y={4} cx={5} cy={6}'.format(
                    s['index'], j, sh['role'], sh['x'], sh['y'], sh['cx'], sh['cy']))
    if bad:
        raise VerifyError('도형이 여백 밖으로 나갔다 — 기하 버그다:\n  ' + '\n  '.join(bad))


# ── 폰트: 이름을 고르고 파일이 있는지만 확인한다 (바이너리 파싱 0줄) ─────────
def norm_name(s):
    return ''.join(ch for ch in str(s).lower() if ch.isalnum() and ch.isascii())


def esc_fc(s):
    out = []
    for ch in str(s):
        if ch in '-:,\\':
            out.append('\\')
        out.append(ch)
    return ''.join(out)


_FC_STATE = {'available': None}


def fc_match(pattern):
    """fc-match 로 이름 하나를 파일 + family 로 해석한다. 없으면 None (한 번 확인하면 기억한다)."""
    if _FC_STATE['available'] is False:
        return None
    try:
        res = subprocess.run(
            ['fc-match', '--format=%{file}|%{family}', pattern],
            stdout=subprocess.PIPE, stderr=subprocess.PIPE,
            universal_newlines=True, timeout=5)
    except FileNotFoundError:
        _FC_STATE['available'] = False
        return None
    except (OSError, subprocess.SubprocessError):
        return None
    if res.returncode != 0 or not res.stdout:
        return None
    _FC_STATE['available'] = True
    cut = res.stdout.find('|')
    if cut < 0:
        return None
    file_path = res.stdout[:cut].strip()
    families = [f.strip() for f in res.stdout[cut + 1:].split(',') if f.strip()]
    if not file_path:
        return None
    return {'file': file_path, 'families': families}


def font_search_dirs():
    home = os.path.expanduser('~')
    if sys.platform == 'darwin':
        return ['/System/Library/Fonts', '/Library/Fonts', os.path.join(home, 'Library', 'Fonts')]
    if sys.platform.startswith('win'):
        dirs = [os.path.join(os.environ.get('WINDIR') or 'C:\\Windows', 'Fonts')]
        local = os.environ.get('LOCALAPPDATA')
        if local:
            dirs.append(os.path.join(local, 'Microsoft', 'Windows', 'Fonts'))
        return dirs
    return [
        '/usr/share/fonts', '/usr/local/share/fonts',
        os.path.join(home, '.local', 'share', 'fonts'), os.path.join(home, '.fonts'),
    ]


def glob_font_file(name):
    """fc-match 가 없는 환경용 파일명 탐색. 깊이·개수 예산을 둬서 병적인 트리에서도 끝난다."""
    hints = FONT_FILE_HINTS.get(str(name).lower(), (norm_name(name),))
    budget = [20000]

    def walk(directory, depth):
        if depth > 6 or budget[0] <= 0:
            return None
        try:
            entries = list(os.scandir(directory))
        except OSError:
            return None
        subdirs = []
        for ent in entries:
            budget[0] -= 1
            if budget[0] <= 0:
                return None
            full = os.path.join(directory, ent.name)
            try:
                if ent.is_dir():
                    subdirs.append(full)
                    continue
            except OSError:
                continue
            stem, ext = os.path.splitext(ent.name)
            if ext.lower() not in FONT_EXT:
                continue
            flat = norm_name(stem)
            if any(h in flat for h in hints):
                return full
        for sub in subdirs:
            hit = walk(sub, depth + 1)
            if hit:
                return hit
        return None

    for directory in font_search_dirs():
        if not os.path.exists(directory):
            continue
        hit = walk(directory, 0)
        if hit:
            return hit
    return None


def find_font_file(name):
    """
    이름 하나를 실제 파일로 해석한다.
    fc-match 는 없는 폰트에도 대체 폰트를 돌려주므로(예: Malgun Gothic → DejaVu Sans)
    돌려받은 family 가 요청한 이름과 실제로 맞는지 대조해야 한다. 안 그러면 이름이 거짓말이 된다.
    """
    want = norm_name(name)
    hit = fc_match(esc_fc(name))
    if hit:
        for fam in hit['families']:
            n = norm_name(fam)
            if n == want or want in n or n in want:
                if os.path.exists(hit['file']):
                    return hit['file']
                break
    globbed = glob_font_file(name)
    if globbed and os.path.exists(globbed):
        return globbed
    return None


def resolve_font(font_arg, meta_font):
    """
    체인을 돌며 "존재하는 파일을 가진 이름" 을 고른다.
    명시된 이름(--font·env·meta.font)이 해석되지 않으면 경고를 남기고 다음으로 내려간다 —
    조용히 바꾸지 않는다.
    """
    warnings = []
    candidates = [
        {'name': font_arg, 'from': '--font'},
        {'name': os.environ.get('BANKER_PPTX_FONT'), 'from': 'BANKER_PPTX_FONT'},
        {'name': meta_font if (meta_font and meta_font != 'auto') else '', 'from': 'meta.font'},
    ]
    explicit = [c for c in candidates if isinstance(c['name'], str) and c['name'].strip() != '']

    for c in explicit:
        name = c['name'].strip()
        file_path = find_font_file(name)
        if file_path:
            return {'name': name, 'file': file_path, 'from': c['from'], 'warnings': warnings}
        warnings.append({'code': 'FONT_NOT_FOUND', 'from': c['from'], 'detail': name})

    for name in FONT_CHAIN:
        file_path = find_font_file(name)
        if file_path:
            return {'name': name, 'file': file_path, 'from': 'chain', 'warnings': warnings}

    sys_hit = fc_match('sans-serif:lang=ko')
    if sys_hit and sys_hit['families'] and os.path.exists(sys_hit['file']):
        return {'name': sys_hit['families'][0], 'file': sys_hit['file'],
                'from': 'system', 'warnings': warnings}

    fallback = explicit[0]['name'].strip() if explicit else FONT_CHAIN[0]
    warnings.append({'code': 'FONT_FILE_UNRESOLVED', 'from': 'chain', 'detail': fallback})
    return {'name': fallback, 'file': None, 'from': 'unresolved', 'warnings': warnings}


# ── python-pptx 해석 (설치는 하지 않는다) ───────────────────────────────────
class Pptx(object):
    """python-pptx 에서 필요한 이름만 모아 둔 묶음. import 를 한 곳에 가둔다."""

    def __init__(self, presentation, emu, pt, rgb, align, anchor, qn):
        self.Presentation = presentation
        self.Emu = emu
        self.Pt = pt
        self.RGBColor = rgb
        self.PP_ALIGN = align
        self.MSO_ANCHOR = anchor
        self.qn = qn


def pptx_missing_message(reason):
    """자동 설치는 하지 않는다. 사람이 읽고 스스로 고칠 수 있는 안내만 낸다(자식 프로세스 0개)."""
    return '\n'.join([
        'python-pptx 를 찾지 못했다. 이 스킬은 의존성을 자동 설치하지 않는다.',
        '',
        '  설치:  python -m pip install python-pptx',
        '         (전역 설치가 막힌 환경이면 --target <dir> 로 받고 PYTHONPATH 에 그 경로를 넣는다)',
        '  안내:  docs-setup 스킬이 이 스킬의 런타임 준비 절차를 담고 있다 (/banker:docs-setup).',
        '  대안:  node 를 쓸 수 있으면 build_a4p.js (pptxgenjs) 로 같은 deck JSON 을 빌드한다.',
        '         검증기 verify_a4p.mjs 는 node 전용이라 검증 경로는 어느 쪽이든 node 를 요구한다.',
        '',
        '  인터프리터: {0} (python {1}.{2}.{3})'.format(
            sys.executable or '(알 수 없음)', *sys.version_info[:3]),
        '  sys.path:   {0}'.format(', '.join(p for p in sys.path if p) or '(비어 있음)'),
        '  원인:       {0}'.format(reason),
    ])


def load_pptx():
    if sys.version_info < MIN_PYTHON:
        raise EnvError(
            'python {0}.{1} 이상이 필요하다 (실행 중인 인터프리터: {2} — python {3}.{4}.{5}).\n'
            '  이 파일은 셰방에 의존하지 않는다. 쓸 인터프리터를 직접 지정해 부른다:\n'
            '    /path/to/python{0}.{1}+ build_a4p.py --in deck.json --out out.pptx\n'
            '  안내: docs-setup 스킬 (/banker:docs-setup).'.format(
                MIN_PYTHON[0], MIN_PYTHON[1], sys.executable or '(알 수 없음)',
                *sys.version_info[:3]))
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.oxml.ns import qn
    except ImportError as exc:
        raise EnvError(pptx_missing_message(exc))
    return Pptx(Presentation, Emu, Pt, RGBColor, PP_ALIGN, MSO_ANCHOR, qn)


# ── 빌드 ────────────────────────────────────────────────────────────────────
def strip_sldsz_type(prs, qn):
    """
    F3. slide_width/height 를 A4 로 바꿔도 python-pptx 는 기본 템플릿의
    type="screen4x3" 선언을 그대로 들고 나간다. 치수는 A4 인데 선언은 4:3 인 파일이고,
    선언 쪽을 믿는 뷰어·변환기가 있다. OOXML 에서 이 속성의 기본값이 custom 이므로
    지우는 것이 곧 custom 선언이다.
    """
    sldsz = prs._element.find(qn('p:sldSz'))
    if sldsz is None:
        raise VerifyError('p:sldSz 를 찾지 못했다 — python-pptx 기본 템플릿이 예상과 다르다.')
    sldsz.attrib.pop('type', None)
    return sldsz


def fix_body_pr(tf, mod):
    """
    F13. a:bodyPr 를 node 빌더와 한 글자까지 맞춘다:
      <a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0" rtlCol="0" anchor="t">
    autofit 요소는 하나도 두지 않는다. python-pptx 의 add_textbox 기본값은
    wrap="none" + <a:spAutoFit/> 라서, 그대로 두면 텍스트가 안 감기고 뷰어가 열 때
    도형 크기를 바꾼다 — 저장된 a:ext 는 맞는데 화면에서 INV-G 가 깨지는 상태다.
    속성을 넣는 순서가 곧 XML 의 속성 순서다. 위 문자열과 같은 순서로 넣는다.
    """
    tf.word_wrap = True                     # wrap="square"
    tf.margin_left = mod.Emu(0)
    tf.margin_top = mod.Emu(0)
    tf.margin_right = mod.Emu(0)
    tf.margin_bottom = mod.Emu(0)
    body_pr = tf._txBody.bodyPr
    body_pr.set('rtlCol', '0')
    tf.vertical_anchor = mod.MSO_ANCHOR.TOP  # anchor="t"
    for tag in ('a:spAutoFit', 'a:normAutofit', 'a:noAutofit'):
        el = body_pr.find(mod.qn(tag))
        if el is not None:
            body_pr.remove(el)


def style_paragraph(paragraph, mod):
    """문단 속성도 node 산출물과 맞춘다: 왼쪽 정렬 · 들여쓰기 0 · 줄간 1.2 · 불릿 글리프 없음."""
    paragraph.alignment = mod.PP_ALIGN.LEFT
    p_pr = paragraph._p.get_or_add_pPr()
    p_pr.set('indent', '0')
    p_pr.set('marL', '0')
    paragraph.line_spacing = LINE_SPACING
    # 불릿 글리프는 항목 텍스트 앞의 '• ' 로 직접 그린다. 목록 자동 글리프가 붙으면 두 번 찍힌다.
    if p_pr.find(mod.qn('a:buNone')) is None:
        p_pr.append(p_pr.makeelement(mod.qn('a:buNone'), {}))


def set_run(run, text, pt, bold, color, typeface, mod):
    """
    F4. run.font.name 은 a:latin 하나만 낸다.
    a:ea 가 없으면 한글이 대체 폰트로 떨어져 tofu 가 나고, 검증기의 폭 판정도 이름을 잘못 읽는다.
    세 슬롯에 같은 typeface 를 직접 넣는다(패턴 출처: skills/arch-diagram/build_pptx_template.py).
    """
    run.text = text
    font = run.font
    font.size = mod.Pt(pt)
    if bold:
        font.bold = True
    font.color.rgb = mod.RGBColor.from_string(color)
    font.name = typeface
    r_pr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = r_pr.find(mod.qn(tag))
        if el is None:
            el = r_pr.makeelement(mod.qn(tag), {})
            r_pr.append(el)
        el.set('typeface', typeface)


# OPC core property `title` 의 길이 상한. OOXML 규격이 아니라 python-pptx 가 스스로 거는 제한인데,
# 두 빌더가 같은 deck 에서 갈리면 원칙 5(교체 가능)가 깨지므로 node 도 똑같이 자른다.
# meta.title 은 문서 메타데이터일 뿐이고 페이지 제목은 별도 도형이라, 잘려도 보이는 내용은 잃지 않는다.
CORE_TITLE_MAX = 255


def clamp_core_title(title):
    if not isinstance(title, str) or len(title) <= CORE_TITLE_MAX:
        return title, False
    return title[:CORE_TITLE_MAX], True


def build(laid, out_path, font, deck, mod, warnings=None):
    prs = mod.Presentation()
    # 치수를 먼저 A4 로 바꾸고, 남는 4:3 선언을 지운다.
    prs.slide_width = mod.Emu(SLIDE_W)
    prs.slide_height = mod.Emu(SLIDE_H)
    strip_sldsz_type(prs, mod.qn)

    meta = deck.get('meta') or {}
    if meta.get('title'):
        value, truncated = clamp_core_title(meta['title'])
        if truncated and warnings is not None:
            warnings.append({'code': 'TITLE_TRUNCATED', 'from': 'meta.title',
                             'detail': '{0} -> {1}'.format(len(meta['title']), CORE_TITLE_MAX)})
        # 클램프가 막아 주지만, 라이브러리가 거는 다른 상한이 남아 있을 수 있다.
        # 원시 트레이스백이 새면 exit code 계약(입력 문제 = 3)이 깨지므로 UsageError 로 들여보낸다.
        try:
            prs.core_properties.title = value
        except ValueError as err:
            raise UsageError('meta.title 을 core property 로 쓸 수 없다: {0}'.format(err))

    # 빈 레이아웃만 쓴다. 기본 템플릿의 자리표시자는 4:3 좌표를 들고 있어서
    # A4 로 바꾼 슬라이드에 얹으면 좌표가 규격 밖으로 나간다.
    blank = prs.slide_layouts[6]

    for s in laid['slides']:
        slide = prs.slides.add_slide(blank)
        for sh in s['shapes']:
            box = slide.shapes.add_textbox(
                mod.Emu(sh['x']), mod.Emu(sh['y']), mod.Emu(sh['cx']), mod.Emu(sh['cy']))
            tf = box.text_frame
            fix_body_pr(tf, mod)
            # 개행은 문단으로 편다 — node 쪽(pptxgenjs)이 '\n' 마다 <a:p> 를 하나씩 내고,
            # 빈 줄에는 run 을 만들지 않는다. 그 모양을 그대로 따른다.
            for i, line in enumerate(str(sh['text']).split('\n')):
                paragraph = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                style_paragraph(paragraph, mod)
                if line != '':
                    set_run(paragraph.add_run(), line, sh['pt'], sh['bold'],
                            sh['color'], font['name'], mod)

    prs.save(out_path)
    return out_path


# ── 출력 ────────────────────────────────────────────────────────────────────
def human_report(laid, font, out_path, warnings):
    g = laid['grid']
    shapes = sum(len(s['shapes']) for s in laid['slides'])
    pages = len(set(s['pageIndex'] for s in laid['slides']))
    lines = [
        'vertical-pptx (python 빌더)',
        '  용지      A4 세로 210x297mm · sldSz {0} x {1} EMU'.format(SLIDE_W, SLIDE_H),
        '  여백      {0}mm · 그리드 {1}열 x {2}행 · 거터 4mm'.format(g['marginMm'], COLS, ROWS),
        '  칸        컬럼 {0} · 행 {1} · 피치 {2} / {3} EMU'.format(
            g['colW'], g['rowH'], g['colPitch'], g['rowPitch']),
    ]
    if font:
        note = '' if font['file'] else '  (파일 미해석 — 검증기 폭 판정은 estimate 로 내려간다)'
        lines.append('  폰트      {0}{1}'.format(font['name'], note))
        lines.append('            {0}'.format(font['file'] or '(경로 없음)'))
    lines.append('  본문      {0}pt · 줄간 {1}'.format(fmt_num(laid['bodyPt']), LINE_SPACING))
    lines.append('  슬라이드  {0} (deck 페이지 {1}개)'.format(len(laid['slides']), pages))
    lines.append('  도형      {0}'.format(shapes))
    lines.append('  경고      {0}'.format(len(warnings)))
    for w in warnings:
        if w.get('page') is not None:
            where = 'page {0} block {1}'.format(w['page'], w['block'])
        elif w.get('from') is not None:
            where = w['from']
        else:
            where = ''
        lines.append('    {0}  {1}  {2}'.format(w['code'], where, w.get('detail') or '').rstrip())
    if out_path:
        lines.append('  출력      {0}'.format(out_path))
    return '\n'.join(lines)


def fmt_num(value):
    """10.5 는 10.5 로, 11.0 은 11 로 — node 의 수 출력과 같게 보이도록."""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


def status_json(font, slides, warnings):
    return json.dumps({
        'skill': 'vertical-pptx',
        'builder': 'python',
        'font': font['name'] if font else None,
        'fontFile': font['file'] if font else None,
        'slides': slides,
        'warnings': warnings,
    }, ensure_ascii=False, separators=(',', ':'))


# ── CLI ─────────────────────────────────────────────────────────────────────
USAGE = '\n'.join([
    'vertical-pptx / build_a4p.py — A4 세로(210x297mm) PPTX python 빌더',
    '',
    '사용법:',
    '  python build_a4p.py --in deck.json --out out.pptx [--font <이름>]',
    '',
    '옵션:',
    '  --in <deck.json>   deck JSON 입력 (두 빌더의 유일한 계약)',
    '  --out <경로>       .pptx 출력 경로',
    '  --font <이름>      폰트 이름을 직접 지정. 체인: --font → BANKER_PPTX_FONT → meta.font',
    '                     → Noto Sans CJK KR → Malgun Gothic → Apple SD Gothic Neo → NanumGothic'
    ' → 시스템 한국어 sans',
    '  -h, --help         이 도움말',
    '',
    '입력은 deck JSON 뿐이다. 마크다운 변환 플래그는 이 빌더에 없다 —',
    '원고를 deck JSON 으로 바꾸는 일은 node 빌더(build_a4p.js)가 단독으로 맡는다.',
    '이 프로젝트의 마크다운 파서는 그 한 벌뿐이고, 두 벌이 되는 순간 같은 원고가 두 개의 덱이 된다.',
    '',
    '요구 사항: python 3.9 이상 · python-pptx (자동 설치하지 않는다)',
    '이 파일은 셰방에 의존하지 않는다. 쓸 인터프리터를 직접 지정해 부르는 것이 안전하다.',
    '',
    '종료 코드: 0 통과 · 1 검증 실패 · 2 환경 미비(python 버전·python-pptx) · 3 사용법 오류',
])


def parse_args(argv):
    opts = {'in': '', 'out': '', 'fontArg': '', 'help': False}
    i = 0
    while i < len(argv):
        a = argv[i]

        def need(name):
            if i + 1 >= len(argv) or argv[i + 1].startswith('--'):
                raise UsageError('{0} 에 값이 필요하다.'.format(name))
            return argv[i + 1]

        if a == '--in':
            opts['in'] = need('--in')
            i += 1
        elif a == '--out':
            opts['out'] = need('--out')
            i += 1
        elif a == '--font':
            opts['fontArg'] = need('--font')
            i += 1
        elif a in ('-h', '--help'):
            opts['help'] = True
        else:
            raise UsageError('알 수 없는 인자: {0}\n\n{1}'.format(a, USAGE))
        i += 1
    return opts


def read_deck_file(path):
    try:
        with open(path, 'r', encoding='utf-8') as fh:
            text = fh.read()
    except OSError as exc:
        raise UsageError('deck JSON 을 읽지 못했다: {0}'.format(exc))
    try:
        return json.loads(text)
    except ValueError as exc:
        raise UsageError('deck JSON 파싱 실패 ({0}): {1}'.format(path, exc))


def run_main(argv):
    opts = parse_args(argv)
    if opts['help'] or not opts['in']:
        sys.stdout.write(USAGE + '\n')
        return EXIT_OK if opts['help'] else EXIT_USAGE
    if not opts['out']:
        raise UsageError('--out <경로> 가 필요하다.')

    deck = read_deck_file(opts['in'])
    laid = layout(deck)  # validate_deck 포함 — 스키마 위반은 여기서 exit 3
    warnings = list(laid['warnings'])
    assert_inside_slide(laid)

    # ★ python-pptx 해석을 폰트 해석보다 먼저 한다.
    #   환경이 미비하면 자식 프로세스(fc-match)를 하나도 띄우지 않고 exit 2 로 나간다.
    mod = load_pptx()
    meta = deck.get('meta') or {}
    font = resolve_font(opts['fontArg'], meta.get('font'))
    warnings.extend(font['warnings'])

    out_path = build(laid, os.path.abspath(opts['out']), font, deck, mod, warnings)
    sys.stdout.write(human_report(laid, font, out_path, warnings) + '\n')
    sys.stdout.write(status_json(font, len(laid['slides']), warnings) + '\n')
    return EXIT_OK


def main(argv):
    try:
        return run_main(argv)
    except (UsageError, EnvError, VerifyError) as err:
        # 환경 미비(exit 2)는 사용자가 그대로 따라 할 수 있는 안내라서 stdout 으로 낸다.
        # 나머지 오류는 stderr.
        stream = sys.stdout if err.exit_code == EXIT_ENV else sys.stderr
        stream.write('{0}\n'.format(err))
        return err.exit_code


if __name__ == '__main__':
    sys.exit(main(sys.argv[1:]))
